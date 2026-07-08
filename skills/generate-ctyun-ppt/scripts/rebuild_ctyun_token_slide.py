from pathlib import Path
import os
import zipfile
from datetime import datetime

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PROJECT = Path(os.environ.get("PPT_PROJECT", Path(__file__).resolve().parents[1]))
TEMPLATE = os.environ.get("PPT_TEMPLATE", "")
MEDIA = PROJECT / "analysis" / "template_media" / "ppt" / "media"
EXPORTS = PROJECT / "exports"
VALIDATION = PROJECT / "validation"
ASSETS = PROJECT / "assets"

W, H = 13.333333, 7.5
RED = RGBColor(190, 15, 23)
DARK_RED = RGBColor(174, 8, 15)
BLACK = RGBColor(31, 31, 31)
GREY = RGBColor(184, 184, 184)
LIGHT_GREY = RGBColor(226, 226, 226)
BLUE = RGBColor(0, 76, 160)


def inch(x):
    return Inches(x)


def add_textbox(slide, x, y, w, h, text="", size=12, color=BLACK, bold=False,
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.02):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = inch(margin)
    tf.margin_right = inch(margin)
    tf.margin_top = inch(margin)
    tf.margin_bottom = inch(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_line(slide, x1, y1, x2, y2, color=LIGHT_GREY, width=1, dash=None):
    line = slide.shapes.add_connector(1, inch(x1), inch(y1), inch(x2), inch(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = dash
    return line


def add_rect(slide, x, y, w, h, fill=None, line=None, radius=False, width=1):
    typ = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(typ, inch(x), inch(y), inch(w), inch(h))
    try:
        shp.shadow.inherit = False
        shp.shadow.visible = False
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(width)
    return shp


def crop_logo():
    ASSETS.mkdir(exist_ok=True)
    src = MEDIA / "image2.png"
    out = ASSETS / "ctyun_logo.png"
    if not src.exists() and TEMPLATE:
        extract_template_media(Path(TEMPLATE))
    if not src.exists():
        raise FileNotFoundError(
            "Template logo source not found. Extract the source PPTX media to "
            f"{MEDIA}, or set PPT_TEMPLATE to the source template PPTX."
        )
    im = Image.open(src).convert("RGBA")
    alpha = im.split()[-1]
    bbox = alpha.getbbox()
    if bbox is None:
        im.save(out)
        return out
    cropped = im.crop(bbox)
    cropped.save(out)
    return out


def extract_template_media(template_path):
    if not template_path.exists():
        raise FileNotFoundError(f"PPT_TEMPLATE not found: {template_path}")
    MEDIA.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(template_path) as archive:
        media_names = [name for name in archive.namelist() if name.startswith("ppt/media/")]
        if not media_names:
            raise FileNotFoundError(f"No ppt/media entries found in template: {template_path}")
        for name in media_names:
            target = PROJECT / "analysis" / "template_media" / name
            target.parent.mkdir(parents=True, exist_ok=True)
            target.write_bytes(archive.read(name))


def add_heading(slide):
    add_rect(slide, 0.30, 0.10, 11.20, 0.78, fill=RGBColor(255, 255, 255), line=GREY, width=1.4)
    slide.shapes[-1].line.dash_style = MSO_LINE_DASH_STYLE.DASH
    add_textbox(slide, 0.45, 0.18, 7.2, 0.50, "南京电信Token运营阶段性进展", 28, RED, True)
    logo = crop_logo()
    slide.shapes.add_picture(str(logo), inch(11.54), inch(0.26), width=inch(1.50))
    add_line(slide, 0.30, 0.96, 13.05, 0.96, RED, 3)
    add_line(slide, 0.30, 0.99, 13.05, 0.99, RGBColor(232, 173, 78), 1.5)


def add_circle_icon(slide, cx, cy, kind):
    circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(cx - 0.42), inch(cy - 0.42), inch(0.84), inch(0.84))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RED
    circ.line.fill.background()
    if kind == "people":
        for dx, dy, r in [(-0.18, -0.08, 0.08), (0, -0.16, 0.09), (0.18, -0.08, 0.08)]:
            head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(cx + dx - r), inch(cy + dy - r), inch(2 * r), inch(2 * r))
            head.fill.solid(); head.fill.fore_color.rgb = RGBColor(255, 255, 255); head.line.fill.background()
        for dx, ww, hh in [(-0.20, 0.14, 0.19), (-0.06, 0.17, 0.25), (0.13, 0.14, 0.19)]:
            add_rect(slide, cx + dx, cy + 0.03, ww, hh, fill=RGBColor(255, 255, 255), radius=True)
    elif kind == "train":
        add_rect(slide, cx - 0.18, cy - 0.22, 0.42, 0.24, fill=None, line=RGBColor(255, 255, 255), width=2)
        add_line(slide, cx - 0.04, cy + 0.18, cx + 0.17, cy - 0.05, RGBColor(255, 255, 255), 2)
        for dx in [-0.25, -0.08, 0.08, 0.25]:
            add_rect(slide, cx + dx - 0.035, cy + 0.18, 0.07, 0.10, fill=RGBColor(255, 255, 255))
    else:
        cup = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.TRAPEZOID, inch(cx - 0.18), inch(cy - 0.17), inch(0.36), inch(0.28))
        cup.fill.solid(); cup.fill.fore_color.rgb = RGBColor(255, 255, 255); cup.line.fill.background()
        add_rect(slide, cx - 0.05, cy + 0.11, 0.10, 0.17, fill=RGBColor(255, 255, 255))
        add_rect(slide, cx - 0.19, cy + 0.27, 0.38, 0.06, fill=RGBColor(255, 255, 255), radius=True)
        star = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.STAR_5_POINT, inch(cx - 0.08), inch(cy - 0.08), inch(0.16), inch(0.16))
        star.fill.solid(); star.fill.fore_color.rgb = RED; star.line.fill.background()


def add_rich_paragraph(slide, x, y, w, title, body_lines, highlights):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(1.45))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = inch(0)
    tf.margin_right = inch(0.02)
    tf.margin_top = inch(0)
    tf.margin_bottom = inch(0)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = RED
    for line in body_lines:
        p = tf.add_paragraph()
        p.space_before = Pt(2)
        p.line_spacing = 1.18
        cursor = 0
        for token, color, bold in highlights:
            idx = line.find(token, cursor)
            if idx < 0:
                continue
            if idx > cursor:
                rr = p.add_run(); rr.text = line[cursor:idx]; rr.font.name = "Microsoft YaHei"; rr.font.size = Pt(11.4); rr.font.color.rgb = BLACK
            rr = p.add_run(); rr.text = token; rr.font.name = "Microsoft YaHei"; rr.font.size = Pt(11.8); rr.font.color.rgb = color; rr.font.bold = bold
            cursor = idx + len(token)
        if cursor < len(line):
            rr = p.add_run(); rr.text = line[cursor:]; rr.font.name = "Microsoft YaHei"; rr.font.size = Pt(11.4); rr.font.color.rgb = BLACK


def add_left_content(slide):
    add_line(slide, 0.88, 2.40, 0.88, 5.34, GREY, 1, MSO_LINE_DASH_STYLE.DASH)
    for y in [2.40, 4.08]:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(0.83), inch(y), inch(0.10), inch(0.10))
        dot.fill.solid(); dot.fill.fore_color.rgb = RED; dot.line.fill.background()
    add_circle_icon(slide, 0.88, 1.68, "people")
    add_circle_icon(slide, 0.88, 3.36, "train")
    add_circle_icon(slide, 0.88, 5.03, "cup")
    add_rich_paragraph(slide, 1.60, 1.26, 4.85, "1. 聚焦四大客群，云图拓客稳步推进", [
        "围绕 DNS 目标、影视动漫、软件开发、国资 OPC 四大核心",
        "客群，全面推进云图专项拓客行动。全区累计梳理目标客户",
        "1876 家，当前已走访 154 家，挖掘有效 AI 算力商机 51 个，",
        "持续夯实后续商机储备。",
    ], [("1876", RED, True), ("154", RED, True), ("51", RED, True)])
    add_line(slide, 1.50, 2.70, 6.30, 2.70, LIGHT_GREY, 1)
    add_rich_paragraph(slide, 1.60, 2.92, 4.85, "2. 强化培训赋能，园区竞赛加速突破", [
        "分层开展一线营销赋能培训，已覆盖玄武、雨花、浦口、江宁、",
        "溧水、大客户中心、秦淮等 7 大营销单元，并同步输出",
        "TokenHub 平台演示、行业解决方案及营销话术材料。玄武、",
        "雨花率先启动首批园区 Token 攻坚 PK 赛，以竞赛牵引园区",
        "精扫走访，持续提升拓客质效。",
    ], [("7", RED, True)])
    add_line(slide, 1.50, 4.56, 6.30, 4.56, LIGHT_GREY, 1)
    add_rich_paragraph(slide, 1.60, 4.78, 4.85, "3. 经营成效显现，标杆项目实现破局", [
        "在本轮 Token 短程攻坚中，成功中标区域头部文化传媒企业",
        "Seedance2.0 视频制作项目，形成南京本地视频类 Token 商用首单，",
        "总签约金额达 480 万元，树立影视动漫行业标准化落地样板，",
        "为同类客户复制推广提供成熟范本。",
    ], [("480", RED, True)])


def add_chart_box(slide, x, y, label):
    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PENTAGON, inch(x), inch(y), inch(2.05), inch(0.35))
    tag.fill.solid(); tag.fill.fore_color.rgb = RED; tag.line.fill.background()
    try:
        tag.shadow.inherit = False
        tag.shadow.visible = False
    except Exception:
        pass
    add_textbox(slide, x + 0.25, y + 0.055, 1.55, 0.20, label, 12.5, RGBColor(255, 255, 255), True, margin=0)
    add_line(slide, x + 2.02, y + 0.18, x + 6.10, y + 0.18, RED, 2)
    add_rect(slide, x, y + 0.35, 6.20, 2.48, fill=RGBColor(255, 255, 255), line=RGBColor(236, 132, 132), radius=False, width=0.8)
    gx, gy, gw, gh = x + 0.55, y + 0.68, 5.50, 1.88
    for i, val in enumerate([100, 80, 60, 40, 20, 0]):
        yy = gy + (100 - val) / 100 * gh
        add_line(slide, gx, yy, gx + gw, yy, RGBColor(225, 225, 225), 0.8, MSO_LINE_DASH_STYLE.DASH)
        add_textbox(slide, gx - 0.38, yy - 0.08, 0.25, 0.14, str(val), 8.5, RGBColor(160, 160, 160), align=PP_ALIGN.RIGHT, margin=0)
    add_line(slide, gx, gy, gx, gy + gh, RGBColor(220, 220, 220), 1)
    add_line(slide, gx, gy + gh, gx + gw, gy + gh, RGBColor(220, 220, 220), 1)
    for i in range(1, 8):
        xx = gx + i * gw / 8
        add_line(slide, xx, gy + gh, xx, gy + gh - 0.04, RGBColor(220, 220, 220), 0.8)
    add_textbox(slide, x + 2.45, y + 1.33, 1.60, 0.30, "图表预留区", 15, RGBColor(180, 180, 180), True, align=PP_ALIGN.CENTER, margin=0)


def add_bottom_cards(slide):
    add_rect(slide, 0.36, 6.20, 5.98, 0.83, fill=RGBColor(255, 255, 255), line=RGBColor(201, 201, 201), radius=True, width=0.8)
    cards = [
        (0.48, "target", "目标客户", "1876", "家"),
        (2.56, "walk", "累计走访", "154", "家"),
        (4.52, "bars", "有效商机", "51", "个"),
    ]
    for idx, (x, kind, label, num, unit) in enumerate(cards):
        add_rect(slide, x, 6.27, 1.88, 0.62, fill=RGBColor(253, 241, 241), radius=True)
        if idx:
            add_line(slide, x - 0.14, 6.24, x - 0.14, 6.98, RGBColor(213, 213, 213), 1)
        add_small_icon(slide, x + 0.42, 6.58, kind)
        add_textbox(slide, x + 0.92, 6.39, 0.88, 0.18, label, 10.5, BLACK, True, align=PP_ALIGN.CENTER, margin=0)
        add_textbox(slide, x + 0.84, 6.61, 0.68, 0.24, num, 16, RED, True, align=PP_ALIGN.RIGHT, margin=0)
        add_textbox(slide, x + 1.54, 6.68, 0.26, 0.14, unit, 8.5, RED, True, align=PP_ALIGN.LEFT, margin=0)


def add_small_icon(slide, cx, cy, kind):
    if kind == "target":
        for r, lw in [(0.19, 3), (0.11, 3)]:
            c = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(cx - r), inch(cy - r), inch(2*r), inch(2*r))
            c.fill.background(); c.line.color.rgb = RED; c.line.width = Pt(lw)
        add_line(slide, cx + 0.02, cy - 0.02, cx + 0.25, cy - 0.26, RED, 3)
    elif kind == "walk":
        head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(cx - 0.05), inch(cy - 0.28), inch(0.10), inch(0.10))
        head.fill.solid(); head.fill.fore_color.rgb = RED; head.line.fill.background()
        add_line(slide, cx, cy - 0.16, cx - 0.03, cy + 0.08, RED, 4)
        add_line(slide, cx - 0.03, cy - 0.04, cx - 0.20, cy + 0.02, RED, 3)
        add_line(slide, cx - 0.02, cy - 0.04, cx + 0.16, cy + 0.05, RED, 3)
        add_line(slide, cx - 0.03, cy + 0.08, cx - 0.17, cy + 0.28, RED, 4)
        add_line(slide, cx - 0.03, cy + 0.08, cx + 0.14, cy + 0.26, RED, 4)
    else:
        for i, h in enumerate([0.16, 0.26, 0.38]):
            add_rect(slide, cx - 0.22 + i * 0.17, cy + 0.20 - h, 0.10, h, fill=RED)
        add_line(slide, cx - 0.25, cy - 0.12, cx - 0.07, cy - 0.26, RED, 3)
        add_line(slide, cx - 0.07, cy - 0.26, cx + 0.15, cy - 0.09, RED, 3)
        add_line(slide, cx + 0.15, cy - 0.09, cx + 0.27, cy - 0.28, RED, 3)


def build():
    EXPORTS.mkdir(exist_ok=True)
    VALIDATION.mkdir(exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    add_heading(slide)
    add_left_content(slide)
    add_line(slide, 6.65, 1.38, 6.65, 7.00, RGBColor(232, 142, 142), 1)
    add_chart_box(slide, 6.78, 1.22, "图表预留区（上）")
    add_chart_box(slide, 6.78, 4.18, "图表预留区（下）")
    add_bottom_cards(slide)
    add_textbox(slide, 12.90, 7.28, 0.18, 0.16, "4", 13, RGBColor(115, 115, 115), align=PP_ALIGN.CENTER, margin=0)

    out = EXPORTS / f"南京电信Token运营阶段性进展_{datetime.now():%Y%m%d_%H%M%S}.pptx"
    for old in EXPORTS.glob("南京电信Token运营阶段性进展_*.pptx"):
        old.unlink()
    prs.save(out)
    (VALIDATION / "readback.md").write_text(
        "\n".join([
            "南京电信Token运营阶段性进展",
            "1. 聚焦四大客群，云图拓客稳步推进",
            "2. 强化培训赋能，园区竞赛加速突破",
            "3. 经营成效显现，标杆项目实现破局",
            "目标客户 1876家",
            "累计走访 154家",
            "有效商机 51个",
            "图表预留区（上）",
            "图表预留区（下）",
            "4",
        ]),
        encoding="utf-8",
    )
    return out


if __name__ == "__main__":
    print(build())
