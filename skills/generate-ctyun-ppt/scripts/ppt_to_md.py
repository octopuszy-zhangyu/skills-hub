#!/usr/bin/env python3
"""
ctyun PPT Replicator - PPTX Text Readback

Extract text from PPTX slides into Markdown for word-for-word validation.

Usage:
    python scripts/ppt_to_md.py <pptx> [-o output.md]

Examples:
    python scripts/ppt_to_md.py exports/deck.pptx -o validation/readback.md

Dependencies:
    python-pptx
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path
from typing import Optional

from pptx import Presentation


def iter_shape_text(shape) -> list[str]:
    """Return non-empty text runs from a shape and its grouped children."""
    lines: list[str] = []
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            lines.extend(iter_shape_text(child))
        return lines

    if not getattr(shape, "has_text_frame", False):
        return lines

    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if text:
            lines.append(text)
    return lines


def pptx_to_markdown(path: Path) -> str:
    """Extract slide text in slide and shape traversal order."""
    prs = Presentation(str(path))
    chunks: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        chunks.append(f"## Slide {index}")
        slide_lines: list[str] = []
        for shape in slide.shapes:
            slide_lines.extend(iter_shape_text(shape))
        chunks.extend(slide_lines or [""])
        chunks.append("")
    return "\n".join(chunks).rstrip() + "\n"


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Extract text from a PPTX into Markdown.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("pptx", help="Input PPTX file.")
    parser.add_argument("-o", "--output", help="Output Markdown path. Defaults to stdout.")
    return parser


def main(argv: Optional[list[str]] = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)

    pptx_path = Path(args.pptx)
    if not pptx_path.exists():
        print(f"PPTX not found: {pptx_path}", file=sys.stderr)
        return 1

    markdown = pptx_to_markdown(pptx_path)
    if args.output:
        output_path = Path(args.output)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(markdown, encoding="utf-8")
    else:
        print(markdown, end="")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
